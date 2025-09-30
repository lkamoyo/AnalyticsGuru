from pptx import Presentation
from pptx.util import Inches
from datetime import datetime

def create_pptx(dashboard_images, insights):
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "AnalyticsGuru Dashboard"
    title_slide.placeholders[1].text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"

    for i, img_path in enumerate(dashboard_images):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(img_path, Inches(1), Inches(1), width=Inches(8))
        textbox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
        textbox.text = insights[i]

    prs.save("dashboard_export.pptx")
    return "dashboard_export.pptx"